"""
generate_docs.py
Generates Technical_Report.docx and Study_Assistant_AI_Presentation.pptx
from the content of Technical_Report.md
"""

import os

# ──────────────────────────────────────────────
# 1. DOCX Generation
# ──────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

doc = Document()

# ── Page margins ──
for section in doc.sections:
    section.top_margin    = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin   = Inches(1.2)
    section.right_margin  = Inches(1.2)

# ── Styles helper ──
styles = doc.styles

def set_heading_style(paragraph, size, color_hex="1F3864", bold=True):
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run(paragraph.text)
    run.bold = bold
    run.font.size = Pt(size)
    r, g, b = int(color_hex[:2], 16), int(color_hex[2:4], 16), int(color_hex[4:], 16)
    run.font.color.rgb = RGBColor(r, g, b)

# ── Title page block ──
title_p = doc.add_paragraph()
title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = title_p.add_run("Study Assistant AI")
run.bold = True
run.font.size = Pt(28)
run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

sub_p = doc.add_paragraph()
sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run2 = sub_p.add_run("Technical Report")
run2.bold = True
run2.font.size = Pt(18)
run2.font.color.rgb = RGBColor(0x7C, 0x3A, 0xED)

doc.add_paragraph()  # spacer

meta_lines = [
    ("Project Title:", "Study Assistant AI"),
    ("Version:", "1.0.0"),
    ("Date:", "March 2026"),
    ("Authors:", "Alok K"),
    ("Technology:", "React · Flask · Groq AI (LLaMA 3.3 70B)"),
]
for label, value in meta_lines:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r1 = p.add_run(label + " ")
    r1.bold = True
    r1.font.size = Pt(11)
    r2 = p.add_run(value)
    r2.font.size = Pt(11)

doc.add_page_break()

# ── Helper: add a section heading ──
def h1(text):
    p = doc.add_heading(level=1)
    p.clear()
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
    return p

def h2(text):
    p = doc.add_heading(level=2)
    p.clear()
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = RGBColor(0x7C, 0x3A, 0xED)
    return p

def body(text):
    p = doc.add_paragraph(text)
    p.style.font.size = Pt(11)
    return p

def bullet(text):
    p = doc.add_paragraph(style="List Bullet")
    run = p.add_run(text)
    run.font.size = Pt(11)
    return p

# ── Table helper ──
def add_table(headers, rows):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.style = "Table Grid"
    # header row
    hdr_cells = t.rows[0].cells
    for i, h in enumerate(headers):
        hdr_cells[i].text = h
        run = hdr_cells[i].paragraphs[0].runs[0]
        run.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # shade header cell
        tc = hdr_cells[i]._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), "1F3864")
        tcPr.append(shd)
    # data rows
    for ri, row in enumerate(rows):
        row_cells = t.rows[ri + 1].cells
        for ci, cell in enumerate(row):
            row_cells[ci].text = cell
            row_cells[ci].paragraphs[0].runs[0].font.size = Pt(10)
    doc.add_paragraph()

# ── TOC heading ──
h1("Table of Contents")
toc_items = [
    "1. Abstract", "2. Introduction", "3. Problem Statement",
    "4. Objectives", "5. Technology Stack", "6. System Architecture",
    "7. Module Descriptions", "8. API Design", "9. AI Integration",
    "10. Frontend Design & UI/UX", "11. Data Management",
    "12. Security Considerations", "13. Testing & Verification",
    "14. Deployment Guide", "15. Future Enhancements",
    "16. Conclusion", "17. References",
]
for item in toc_items:
    bullet(item)
doc.add_page_break()

# ══════════════════ SECTIONS ══════════════════

# 1. Abstract
h1("1. Abstract")
body(
    "Study Assistant AI is a full-stack web application designed to help students learn more "
    "effectively using artificial intelligence. The application provides four core study tools — "
    "Notes Management, AI-powered Flashcard Generation, Dynamic Quiz Creation, and an AI Chat "
    "Tutor — all integrated into a single, modern, dark-themed interface. The backend is built "
    "with Python Flask and integrates with the Groq cloud AI platform (LLaMA 3.3 70B model) for "
    "intelligent content generation. The frontend is a React single-page application styled with "
    "Ant Design and features glassmorphism aesthetics, animated backgrounds, and a fully "
    "responsive layout. Data persistence is handled through JSON file storage, making the system "
    "lightweight and easy to deploy without an external database."
)
doc.add_paragraph()

# 2. Introduction
h1("2. Introduction")
body(
    "The rapid evolution of AI technologies has opened new possibilities in the education sector. "
    "Modern students need tools that go beyond simple note-taking — they need intelligent systems "
    "that can generate study material, quiz them on topics, and provide instant explanations."
)
body("Study Assistant AI addresses this need by combining traditional study management features with AI-powered capabilities. The application is designed to be:")
for item in [
    "Accessible — Runs in any modern web browser.",
    "Intelligent — Leverages a 70-billion-parameter language model for high-quality content generation.",
    "Beautiful — Features a premium dark-mode UI with glassmorphism, animations, and micro-interactions.",
    "Lightweight — No external database required; data is stored as JSON files.",
]:
    bullet(item)
doc.add_paragraph()

# 3. Problem Statement
h1("3. Problem Statement")
body("Students often struggle with:")
for item in [
    "Disorganized study materials — Notes scattered across multiple platforms.",
    "Passive learning — Simply reading notes without active recall practice.",
    "Lack of instant feedback — No way to quickly test understanding of a topic.",
    "Limited access to tutoring — Human tutors are expensive and not always available.",
]:
    bullet(item)
body(
    "Existing solutions either focus on a single feature or require expensive subscriptions. "
    "There is a need for an integrated, AI-powered, free-to-use study platform that combines "
    "all major study tools in one place."
)
doc.add_paragraph()

# 4. Objectives
h1("4. Objectives")
add_table(
    ["#", "Objective", "Status"],
    [
        ["1", "Build a notes management system with CRUD operations", "✅ Completed"],
        ["2", "Implement flashcard creation (manual + AI-generated)", "✅ Completed"],
        ["3", "Create AI-powered quiz generation with auto-grading", "✅ Completed"],
        ["4", "Integrate an AI chat tutor for Q&A assistance", "✅ Completed"],
        ["5", "Design a premium, responsive UI with dark mode", "✅ Completed"],
        ["6", "Ensure lightweight deployment (no external DB)", "✅ Completed"],
    ],
)

# 5. Technology Stack
h1("5. Technology Stack")
h2("5.1 Frontend")
add_table(
    ["Technology", "Version", "Purpose"],
    [
        ["React", "19.2.4", "UI component framework"],
        ["Vite", "8.0.1", "Build tool & dev server"],
        ["Ant Design", "6.3.3", "UI component library"],
        ["React Router DOM", "7.13.1", "Client-side routing"],
        ["Axios", "1.13.6", "HTTP client for API calls"],
        ["@ant-design/icons", "6.1.0", "Icon library"],
    ],
)
h2("5.2 Backend")
add_table(
    ["Technology", "Version", "Purpose"],
    [
        ["Python", "3.8+", "Server-side language"],
        ["Flask", "3.1.0", "Web framework"],
        ["Flask-CORS", "5.0.1", "Cross-Origin Resource Sharing"],
        ["OpenAI SDK", "≥1.0.0", "AI API client (Groq-compatible)"],
        ["python-dotenv", "≥1.0.0", "Environment variable management"],
    ],
)
h2("5.3 AI Platform")
add_table(
    ["Component", "Detail"],
    [
        ["Provider", "Groq Cloud"],
        ["Model", "LLaMA 3.3 70B Versatile"],
        ["API Compatibility", "OpenAI-compatible REST API"],
        ["Tier", "Free tier"],
    ],
)

# 6. System Architecture
h1("6. System Architecture")
h2("6.1 High-Level Architecture")
body(
    "The application follows a three-tier architecture: the React frontend communicates with "
    "the Flask backend via RESTful HTTP/JSON (Axios), and the Flask backend calls the Groq AI "
    "API over HTTPS using the OpenAI-compatible SDK. Data is persisted in JSON files on the server."
)
h2("6.2 Design Pattern")
body("The application follows a layered architecture:")
for layer in [
    "Presentation Layer — React components & pages.",
    "API Layer — Flask route blueprints (thin controllers).",
    "Service Layer — Business logic, data transformation, AI orchestration.",
    "Data Layer — JSON file read/write utilities.",
    "External Layer — Groq AI API for LLM capabilities.",
]:
    bullet(layer)
h2("6.3 Communication Protocol")
for item in [
    "Frontend ↔ Backend: RESTful HTTP/JSON via Axios.",
    "Backend ↔ AI: OpenAI-compatible SDK over HTTPS to Groq cloud.",
    "CORS: Configured to allow requests from http://localhost:5173.",
]:
    bullet(item)
doc.add_paragraph()

# 7. Module Descriptions
h1("7. Module Descriptions")
modules = [
    ("7.1 Notes Module", "Full CRUD management of study notes organized by subject. The backend uses 5 REST endpoints (GET all, GET by ID, POST, PUT, DELETE) with UUID generation and timestamp management. The frontend provides a search-enabled management UI."),
    ("7.2 Flashcards Module", "Manual and AI-powered flashcard creation for active recall. The AI generation endpoint accepts a text block, sends it to the LLM with a structured prompt, and parses the JSON response into {question, answer} pairs. The frontend features a 3D flip card animation."),
    ("7.3 Quiz Module", "AI-generated topic-based quizzes with auto-grading and score history. Each question has 4 multiple-choice options with a 0-based correct_answer index. The frontend provides a progress-tracked quiz interface."),
    ("7.4 Chat Module", "Conversational AI tutor that answers student questions in real-time. The system prompt configures the AI to be smart, friendly, and encouraging, responding in markdown and in the user's language."),
]
for title, desc in modules:
    h2(title)
    body(desc)
doc.add_paragraph()

# 8. API Design
h1("8. API Design")
h2("8.1 Endpoint Summary")
add_table(
    ["Method", "Endpoint", "Description", "Request Body"],
    [
        ["GET",    "/api/notes/",               "Get all notes",            "—"],
        ["POST",   "/api/notes/",               "Create a note",            "{title, content, subject?}"],
        ["PUT",    "/api/notes/<id>",            "Update a note",            "{title?, content?, subject?}"],
        ["DELETE", "/api/notes/<id>",            "Delete a note",            "—"],
        ["GET",    "/api/flashcards/",           "Get all flashcards",       "—"],
        ["POST",   "/api/flashcards/",           "Create a flashcard",       "{question, answer, subject?}"],
        ["POST",   "/api/flashcards/generate",   "AI-generate flashcards",   "{text, subject?, num_cards?}"],
        ["DELETE", "/api/flashcards/<id>",       "Delete a flashcard",       "—"],
        ["POST",   "/api/quiz/generate",         "Generate a quiz",          "{topic, num_questions?}"],
        ["POST",   "/api/quiz/grade",            "Grade a quiz",             "{quiz_id, answers}"],
        ["GET",    "/api/quiz/history",          "Get quiz history",         "—"],
        ["POST",   "/api/chat/",                 "Chat with AI",             "{message}"],
        ["GET",    "/api/health",                "Health check",             "—"],
    ],
)
h2("8.2 Status Codes")
add_table(
    ["Code", "Meaning"],
    [
        ["200", "Success"],
        ["201", "Resource created"],
        ["400", "Bad request / validation error"],
        ["404", "Resource not found"],
        ["500", "Internal server error"],
    ],
)

# 9. AI Integration
h1("9. AI Integration")
h2("9.1 Architecture")
body(
    "The AI integration is centralized in app/services/ai_service.py. It uses the OpenAI "
    "Python SDK pointed at the Groq base URL with the llama-3.3-70b-versatile model."
)
h2("9.2 AI Functions")
add_table(
    ["Function", "Purpose", "Temperature"],
    [
        ["generate_response(message)",           "Chat tutor responses",          "0.7 (balanced)"],
        ["generate_flashcards_from_text(text,n)","Flashcard Q&A pairs",           "0.5 (focused)"],
        ["generate_quiz_questions(topic,n)",     "MCQ generation",                "0.6 (moderate)"],
    ],
)
h2("9.3 Prompt Engineering")
for item in [
    "Chat: Study assistant persona — friendly, encouraging, concise, markdown-formatted.",
    "Flashcards: JSON-only output instruction — generates [{question, answer}] arrays.",
    "Quiz: Structured MCQ generation — 4 options per question with 0-based correct answer index.",
]:
    bullet(item)
h2("9.4 Response Parsing")
body("AI responses are parsed with defensive coding: strip whitespace → remove markdown code fences → parse via json.loads().")
doc.add_paragraph()

# 10. Frontend Design
h1("10. Frontend Design & UI/UX")
h2("10.1 Design Philosophy")
body("The frontend follows a premium dark-mode glassmorphism design language:")
for item in [
    "Color Palette: Deep indigo base (#12122a) with vibrant purple primary (#7c3aed).",
    "Typography: Inter font family with system font fallbacks.",
    "Glass Effects: Semi-transparent containers with backdrop blur.",
    "Animations: Animated gradient background, smooth page transitions, and micro-interactions.",
]:
    bullet(item)
h2("10.2 Routing")
add_table(
    ["Path", "Component", "Description"],
    [
        ["/",           "Dashboard",     "Home page with feature overview"],
        ["/notes",      "NotesPage",     "Notes management"],
        ["/flashcards", "FlashcardsPage","Flashcard viewer"],
        ["/quiz",       "QuizPage",      "Quiz generation & taking"],
        ["/chat",       "ChatPage",      "AI chat interface"],
    ],
)

# 11. Data Management
h1("11. Data Management")
h2("11.1 Storage Strategy")
body("The application uses JSON file-based storage instead of a traditional database:")
add_table(
    ["File", "Contents"],
    [
        ["backend/data/notes.json",     "Array of note objects"],
        ["backend/data/flashcards.json","Array of flashcard objects"],
        ["backend/data/quizzes.json",   "Array of quiz objects with history"],
    ],
)
h2("11.2 Advantages of JSON Storage")
for item in [
    "Zero Setup: No database installation or configuration needed.",
    "Portability: Data files can be easily shared, backed up, or versioned.",
    "Simplicity: Perfect for single-user, educational applications.",
    "Human-Readable: Data can be inspected and edited directly.",
]:
    bullet(item)
doc.add_paragraph()

# 12. Security
h1("12. Security Considerations")
add_table(
    ["Concern", "Implementation"],
    [
        ["API Key Protection",   "Keys stored in .env file (gitignored)"],
        ["CORS",                 "Restricted to frontend origin (localhost:5173)"],
        ["Input Validation",     "All endpoints validate request bodies via schemas.py"],
        ["Error Handling",       "Structured error responses with appropriate HTTP codes"],
        ["Environment Isolation","Python virtual environment (venv) for dependencies"],
    ],
)
h2("12.1 Recommendations for Production")
for item in [
    "Replace SECRET_KEY with a cryptographically secure value.",
    "Use HTTPS in production.",
    "Add rate limiting to AI endpoints to prevent API key abuse.",
    "Implement user authentication (JWT / OAuth 2.0).",
    "Migrate from JSON files to a proper database (PostgreSQL / MongoDB).",
]:
    bullet(item)
doc.add_paragraph()

# 13. Testing
h1("13. Testing & Verification")
h2("13.1 Manual Testing")
body("All features have been manually verified:")
for item in [
    "Notes: Create, read, update, delete operations.",
    "Flashcards: Manual creation and AI generation from text.",
    "Quiz: Topic-based generation, question answering, and grading.",
    "Chat: Multi-turn conversations with AI tutor.",
    "UI/UX: Dark theme, animations, responsive layout across screen sizes.",
]:
    bullet(item)
h2("13.2 API Testing")
body("Endpoints can be tested using curl / Postman for direct API testing, or Browser DevTools for frontend-backend communication inspection.")
h2("13.3 Health Check")
body("GET /api/health → Response: {\"status\": \"ok\", \"message\": \"Study Assistant API is running!\"}")
doc.add_paragraph()

# 14. Deployment Guide
h1("14. Deployment Guide")
h2("14.1 Prerequisites")
for item in [
    "Python 3.8+ installed",
    "Node.js 18+ and npm installed",
    "A Groq API key (free at https://console.groq.com)",
]:
    bullet(item)
h2("14.2 Backend Setup")
body("1. cd backend  2. python -m venv venv  3. venv\\Scripts\\activate (Windows)  4. pip install -r requirements.txt  5. Create .env with GROQ_API_KEY=your_key_here  6. python run.py")
h2("14.3 Frontend Setup")
body("1. cd frontend  2. npm install  3. npm run dev  4. Open http://localhost:5173")
doc.add_paragraph()

# 15. Future Enhancements
h1("15. Future Enhancements")
add_table(
    ["Priority", "Feature", "Description"],
    [
        ["High",   "User Authentication",  "JWT-based login system for multi-user support"],
        ["High",   "Database Migration",   "Replace JSON files with PostgreSQL or MongoDB"],
        ["Medium", "File Upload & OCR",    "Extract text from PDF/images for AI processing"],
        ["Medium", "Spaced Repetition",    "Implement SM-2 algorithm for flashcard scheduling"],
        ["Medium", "Export & Share",       "PDF export for notes, shareable quiz links"],
        ["Low",    "Voice Input",          "Speech-to-text for chat and note dictation"],
        ["Low",    "Mobile App",           "React Native companion app"],
        ["Low",    "Analytics Dashboard",  "Study time tracking and performance trends"],
    ],
)

# 16. Conclusion
h1("16. Conclusion")
body(
    "Study Assistant AI successfully demonstrates how modern AI can be integrated into educational "
    "tools to create a powerful, personalized learning experience. The application combines the "
    "reliability of a traditional full-stack web architecture (React + Flask) with the intelligence "
    "of a state-of-the-art language model (LLaMA 3.3 70B via Groq).\n\n"
    "The modular architecture ensures maintainability and extensibility — new features can be added "
    "by simply creating new route blueprints and service modules. The JSON-based storage makes the "
    "system immediately deployable without any database configuration, while the premium "
    "glassmorphism UI provides an engaging and modern user experience.\n\n"
    "The project serves as both a practical study tool and a reference implementation for building "
    "AI-powered web applications with Python and JavaScript."
)
doc.add_paragraph()

# 17. References
h1("17. References")
refs = [
    "Flask Documentation — https://flask.palletsprojects.com/",
    "React Documentation — https://react.dev/",
    "Ant Design — https://ant.design/",
    "Vite — https://vite.dev/",
    "Groq AI Platform — https://console.groq.com/docs",
    "OpenAI Python SDK — https://github.com/openai/openai-python",
    "LLaMA Model Family — Meta AI Research",
    "React Router — https://reactrouter.com/",
]
for i, ref in enumerate(refs, 1):
    bullet(f"{i}. {ref}")

# ── Save DOCX ──
docx_path = os.path.join(os.path.dirname(__file__), "Technical_Report.docx")
doc.save(docx_path)
print(f"✅ DOCX saved: {docx_path}")

# ──────────────────────────────────────────────
# 2. PPTX Generation (14 slides)
# ──────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color palette ──
C_BG      = PptxRGB(0x12, 0x12, 0x2A)   # deep indigo
C_ACCENT  = PptxRGB(0x7C, 0x3A, 0xED)   # vibrant purple
C_ACCENT2 = PptxRGB(0xA7, 0x8B, 0xFA)   # light purple
C_WHITE   = PptxRGB(0xFF, 0xFF, 0xFF)
C_LGRAY   = PptxRGB(0xCC, 0xCC, 0xEE)
C_GOLD    = PptxRGB(0xF5, 0xD0, 0x60)

BLANK_LAYOUT = prs.slide_layouts[6]  # completely blank

def add_bg(slide, color=C_BG):
    """Fill slide background."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h, font_size=18, bold=False,
                color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_size=14, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def title_slide_base(title, subtitle=None, tag=None):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    # accent bar left
    add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT)
    # bottom bar
    add_rect(slide, 0, 6.8, 13.33, 0.7, C_ACCENT)
    # decorative circle top-right
    circ = slide.shapes.add_shape(9, Inches(10.5), Inches(-1.0), Inches(3.5), Inches(3.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = PptxRGB(0x3D, 0x1A, 0x78)
    circ.line.fill.background()
    # title
    add_textbox(slide, title, 0.4, 2.5, 12.0, 1.4, font_size=40, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    # purple underline
    add_rect(slide, 4.5, 4.0, 4.33, 0.05, C_ACCENT)
    if subtitle:
        add_textbox(slide, subtitle, 0.4, 4.1, 12.0, 0.8, font_size=20,
                    color=C_ACCENT2, align=PP_ALIGN.CENTER)
    if tag:
        add_textbox(slide, tag, 0.3, 6.85, 12.0, 0.55, font_size=13,
                    color=C_WHITE, align=PP_ALIGN.CENTER)
    return slide

def content_slide(title):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    # top accent bar
    add_rect(slide, 0, 0, 13.33, 0.08, C_ACCENT)
    # title area
    add_rect(slide, 0, 0.08, 13.33, 1.1, PptxRGB(0x1A, 0x1A, 0x3E))
    add_textbox(slide, title, 0.3, 0.1, 12.5, 1.0, font_size=26, bold=True,
                color=C_WHITE, align=PP_ALIGN.LEFT)
    # purple left accent
    add_rect(slide, 0, 1.18, 0.06, 6.32, C_ACCENT)
    # footer
    add_rect(slide, 0, 7.0, 13.33, 0.5, PptxRGB(0x1A, 0x1A, 0x3E))
    add_textbox(slide, "Study Assistant AI  |  Technical Presentation  |  March 2026",
                0.3, 7.05, 12.0, 0.4, font_size=10, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    return slide

def add_bullet_box(slide, items, l, t, w, h, font_size=15, bullet_char="▸"):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = C_LGRAY
        p.space_before = Pt(4)
    return txBox

def add_table_ppt(slide, headers, rows, l, t, w):
    cols = len(headers)
    total_rows = 1 + len(rows)
    tbl = slide.shapes.add_table(total_rows, cols, Inches(l), Inches(t),
                                  Inches(w), Inches(0.4 * total_rows)).table
    # header
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = C_ACCENT
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(11)
        p.runs[0].font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
    # rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            bg = PptxRGB(0x1E, 0x1E, 0x44) if ri % 2 == 0 else PptxRGB(0x16, 0x16, 0x36)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(10)
            p.runs[0].font.color.rgb = C_LGRAY
            p.alignment = PP_ALIGN.LEFT
    return tbl

# ════════════ SLIDE 1 — Title ════════════
slide1 = title_slide_base(
    "Study Assistant AI",
    "Technical Report  •  Version 1.0.0",
    "Alok K  |  React · Flask · Groq AI (LLaMA 3.3 70B)  |  March 2026"
)

# ════════════ SLIDE 2 — Agenda ════════════
slide2 = content_slide("Agenda")
agenda = [
    "01  Project Overview & Abstract",
    "02  Problem Statement",
    "03  Objectives",
    "04  Technology Stack",
    "05  System Architecture",
    "06  Module Descriptions",
    "07  API Design",
    "08  AI Integration & Prompt Engineering",
    "09  Frontend Design & UI/UX",
    "10  Data Management",
    "11  Security Considerations",
    "12  Testing & Verification",
    "13  Deployment Guide",
    "14  Future Enhancements & Conclusion",
]
# two columns
left_items  = agenda[:7]
right_items = agenda[7:]
add_bullet_box(slide2, left_items,  0.4, 1.4, 6.2, 5.4, font_size=14, bullet_char="")
add_bullet_box(slide2, right_items, 6.8, 1.4, 6.2, 5.4, font_size=14, bullet_char="")

# ════════════ SLIDE 3 — Abstract / Overview ════════════
slide3 = content_slide("Project Overview")
add_textbox(slide3,
    "Study Assistant AI is a full-stack web application that integrates four intelligent "
    "study tools into a single, modern interface — empowering students to study smarter, "
    "not harder.",
    0.4, 1.3, 12.5, 1.0, font_size=16, color=C_LGRAY)

cards = [
    ("📝", "Notes\nManagement",   "Full CRUD with subject organization"),
    ("🃏", "AI Flashcards",       "Auto-generate Q&A from any text"),
    ("🧩", "Dynamic Quizzes",     "Topic-based MCQs with auto-grading"),
    ("💬", "AI Chat Tutor",       "Real-time conversational learning"),
]
for i, (icon, name, desc) in enumerate(cards):
    x = 0.4 + i * 3.2
    add_rect(slide3, x, 2.55, 3.0, 3.8, PptxRGB(0x1E, 0x1E, 0x4A))
    add_textbox(slide3, icon,  x + 0.1, 2.7,  2.8, 0.6, font_size=24, align=PP_ALIGN.CENTER)
    add_textbox(slide3, name,  x + 0.1, 3.3,  2.8, 0.7, font_size=14, bold=True,
                color=C_ACCENT2, align=PP_ALIGN.CENTER)
    add_textbox(slide3, desc,  x + 0.1, 4.1,  2.8, 1.5, font_size=12,
                color=C_LGRAY, align=PP_ALIGN.CENTER)
    # accent top border on card
    add_rect(slide3, x, 2.55, 3.0, 0.05, C_ACCENT)

# ════════════ SLIDE 4 — Problem Statement ════════════
slide4 = content_slide("Problem Statement")
add_textbox(slide4, "Why does this application exist?", 0.4, 1.3, 12.5, 0.5,
            font_size=16, bold=True, color=C_ACCENT2)
problems = [
    ("📚  Disorganized Materials",  "Notes scattered across multiple platforms with no central hub."),
    ("😴  Passive Learning",         "Students read notes but never test active recall."),
    ("❓  No Instant Feedback",      "No quick way to verify understanding of a topic."),
    ("💰  Limited Tutoring Access",  "Human tutors are expensive and not always available."),
]
for i, (title, desc) in enumerate(problems):
    y = 2.0 + i * 1.1
    add_rect(slide4, 0.3, y, 12.5, 0.95, PptxRGB(0x1E, 0x1E, 0x44))
    add_textbox(slide4, title, 0.5, y + 0.05, 4.5, 0.45, font_size=14, bold=True, color=C_GOLD)
    add_textbox(slide4, desc,  5.2, y + 0.05, 7.5, 0.45, font_size=13, color=C_LGRAY)

# ════════════ SLIDE 5 — Objectives ════════════
slide5 = content_slide("Objectives")
add_table_ppt(slide5,
    ["#", "Objective", "Status"],
    [
        ["1", "Build a notes management system with CRUD operations", "✅ Completed"],
        ["2", "Implement flashcard creation (manual + AI-generated)", "✅ Completed"],
        ["3", "Create AI-powered quiz generation with auto-grading",  "✅ Completed"],
        ["4", "Integrate an AI chat tutor for Q&A assistance",        "✅ Completed"],
        ["5", "Design a premium, responsive UI with dark mode",       "✅ Completed"],
        ["6", "Ensure lightweight deployment (no external DB)",       "✅ Completed"],
    ],
    0.4, 1.5, 12.5
)

# ════════════ SLIDE 6 — Technology Stack ════════════
slide6 = content_slide("Technology Stack")
# Frontend
add_textbox(slide6, "Frontend", 0.4, 1.3, 4.0, 0.4, font_size=15, bold=True, color=C_ACCENT2)
fe = [("React 19","UI Framework"), ("Vite 8","Build Tool"), ("Ant Design 6","UI Library"),
      ("React Router 7","Routing"), ("Axios 1.13","HTTP Client")]
for i, (tech, purpose) in enumerate(fe):
    y = 1.8 + i * 0.8
    add_rect(slide6, 0.3, y, 5.8, 0.68, PptxRGB(0x1E, 0x1E, 0x44))
    add_textbox(slide6, tech,    0.5, y+0.1, 2.5, 0.45, font_size=12, bold=True, color=C_WHITE)
    add_textbox(slide6, purpose, 3.0, y+0.1, 3.0, 0.45, font_size=12, color=C_LGRAY)

# Backend
add_textbox(slide6, "Backend", 6.7, 1.3, 4.0, 0.4, font_size=15, bold=True, color=C_ACCENT2)
be = [("Python 3.8+","Language"), ("Flask 3.1","Web Framework"),
      ("Flask-CORS","CORS Handling"), ("OpenAI SDK","AI Client"), ("python-dotenv","Env Vars")]
for i, (tech, purpose) in enumerate(be):
    y = 1.8 + i * 0.8
    add_rect(slide6, 6.6, y, 6.3, 0.68, PptxRGB(0x1E, 0x1E, 0x44))
    add_textbox(slide6, tech,    6.8, y+0.1, 2.5, 0.45, font_size=12, bold=True, color=C_WHITE)
    add_textbox(slide6, purpose, 9.3, y+0.1, 3.5, 0.45, font_size=12, color=C_LGRAY)

# AI tag at bottom
add_rect(slide6, 3.5, 6.5, 6.3, 0.42, C_ACCENT)
add_textbox(slide6, "🤖  AI: Groq Cloud  •  LLaMA 3.3 70B Versatile  •  Free Tier",
            3.6, 6.52, 6.1, 0.38, font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ════════════ SLIDE 7 — System Architecture ════════════
slide7 = content_slide("System Architecture")
add_textbox(slide7,
    "The application follows a 5-layer architecture with clear separation of concerns.",
    0.4, 1.3, 12.5, 0.5, font_size=14, color=C_LGRAY)
layers = [
    ("Presentation", "React + Vite SPA  (Port 5173)"),
    ("API",          "Flask Route Blueprints  (Port 5000)"),
    ("Service",      "Business Logic, AI Orchestration"),
    ("Data",         "JSON File Read/Write Utilities"),
    ("External",     "Groq AI API  (HTTPS / OpenAI SDK)"),
]
arrow_y = [1.95, 2.75, 3.55, 4.35, 5.15]
for i, ((layer, desc), y) in enumerate(zip(layers, arrow_y)):
    clr = [C_ACCENT, PptxRGB(0x5B, 0x21, 0xB6), PptxRGB(0x4C, 0x1D, 0x95),
           PptxRGB(0x3B, 0x16, 0x78), PptxRGB(0x2E, 0x10, 0x65)][i]
    add_rect(slide7, 0.5, y, 12.0, 0.62, clr)
    add_textbox(slide7, layer, 0.7, y+0.08, 3.0, 0.45, font_size=13, bold=True, color=C_WHITE)
    add_textbox(slide7, desc,  4.0, y+0.08, 8.5, 0.45, font_size=13, color=C_WHITE)
    if i < 4:
        add_textbox(slide7, "↕", 6.0, y + 0.62, 1.0, 0.3, font_size=12, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# ════════════ SLIDE 8 — Module Descriptions ════════════
slide8 = content_slide("Module Descriptions")
mods = [
    ("📝 Notes",        "5 REST endpoints · UUID + timestamp management · Search-enabled CRUD UI"),
    ("🃏 Flashcards",   "AI text→Q&A generation · 3D flip card animation · JSON persistence"),
    ("🧩 Quiz",         "AI MCQ generation · 4-option format · Auto-grading · Score history"),
    ("💬 Chat Tutor",   "OpenAI-compatible chat · Study persona system prompt · Markdown output"),
]
for i, (mod, desc) in enumerate(mods):
    y = 1.5 + i * 1.3
    add_rect(slide8, 0.3, y, 12.5, 1.15, PptxRGB(0x1E, 0x1E, 0x44))
    add_rect(slide8, 0.3, y, 0.35, 1.15, C_ACCENT)
    add_textbox(slide8, mod,  0.8, y+0.1,  3.5, 0.45, font_size=16, bold=True, color=C_ACCENT2)
    add_textbox(slide8, desc, 0.8, y+0.6, 11.5, 0.45, font_size=13, color=C_LGRAY)

# ════════════ SLIDE 9 — API Design ════════════
slide9 = content_slide("API Design")
add_table_ppt(slide9,
    ["Method", "Endpoint", "Description"],
    [
        ["GET",    "/api/notes/",              "Get all notes"],
        ["POST",   "/api/notes/",              "Create a note"],
        ["POST",   "/api/flashcards/generate", "AI-generate flashcards"],
        ["POST",   "/api/quiz/generate",       "Generate a quiz"],
        ["POST",   "/api/quiz/grade",          "Grade a quiz"],
        ["POST",   "/api/chat/",               "Chat with AI tutor"],
        ["GET",    "/api/health",              "Health check"],
    ],
    0.4, 1.5, 12.5
)
add_textbox(slide9,
    "All endpoints return JSON  ·  Errors: { \"error\": \"description\" }  ·  Status codes: 200, 201, 400, 404, 500",
    0.4, 6.3, 12.5, 0.55, font_size=12, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# ════════════ SLIDE 10 — AI Integration ════════════
slide10 = content_slide("AI Integration & Prompt Engineering")
add_textbox(slide10, "Provider: Groq Cloud  •  Model: LLaMA 3.3 70B Versatile  •  OpenAI-compatible SDK",
            0.4, 1.3, 12.5, 0.5, font_size=14, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

funcs = [
    ("generate_response()",              "Chat tutoring",         "0.7", "Friendly study persona, markdown output"),
    ("generate_flashcards_from_text()",  "Flashcard Q&A pairs",   "0.5", "JSON array [{question, answer}]"),
    ("generate_quiz_questions()",        "MCQ generation",        "0.6", "4 options, 0-based correct answer index"),
]
add_table_ppt(slide10,
    ["Function", "Purpose", "Temp", "Output Format"],
    [(a, b, c, d) for a, b, c, d in funcs],
    0.4, 2.0, 12.5
)
add_textbox(slide10, "Response Parsing: strip whitespace → remove ```json fences → json.loads()",
            0.4, 4.5, 12.5, 0.5, font_size=13, color=C_LGRAY, align=PP_ALIGN.CENTER)

# ════════════ SLIDE 11 — Frontend Design ════════════
slide11 = content_slide("Frontend Design & UI/UX")
design_items = [
    ("🎨 Color Palette",  "#12122a deep indigo base  +  #7c3aed vibrant purple primary"),
    ("✍  Typography",     "Inter font family with system font fallbacks"),
    ("🪟 Glass Effects",   "Semi-transparent containers with CSS backdrop-filter blur"),
    ("✨ Animations",      "Animated gradient background, smooth transitions, micro-interactions"),
    ("📱 Responsive",      "Fully responsive layout across all screen sizes"),
]
for i, (label, desc) in enumerate(design_items):
    y = 1.5 + i * 1.0
    add_rect(slide11, 0.3, y, 12.5, 0.85, PptxRGB(0x1E, 0x1E, 0x44))
    add_rect(slide11, 0.3, y, 3.8,  0.85, PptxRGB(0x2D, 0x1A, 0x6E))
    add_textbox(slide11, label, 0.5, y+0.15, 3.5, 0.55, font_size=13, bold=True, color=C_ACCENT2)
    add_textbox(slide11, desc,  4.3, y+0.15, 8.3, 0.55, font_size=13, color=C_LGRAY)

# ════════════ SLIDE 12 — Data & Security ════════════
slide12 = content_slide("Data Management & Security")
# Data section
add_textbox(slide12, "Data Storage", 0.4, 1.3, 5.5, 0.4, font_size=15, bold=True, color=C_ACCENT2)
data_items = [
    "JSON file-based storage (no external DB)",
    "notes.json  ·  flashcards.json  ·  quizzes.json",
    "UUID identifiers + ISO-8601 timestamps",
    "Zero setup — instantly portable & human-readable",
]
add_bullet_box(slide12, data_items, 0.4, 1.8, 5.8, 3.5, font_size=13)

# Security section
add_textbox(slide12, "Security", 7.0, 1.3, 5.5, 0.4, font_size=15, bold=True, color=C_ACCENT2)
sec_items = [
    "API keys in .env (gitignored)",
    "CORS restricted to localhost:5173",
    "Input validation via schemas.py",
    "Structured HTTP error responses",
    "Python venv dependency isolation",
]
add_bullet_box(slide12, sec_items, 7.0, 1.8, 5.9, 3.5, font_size=13)

# divider
add_rect(slide12, 6.6, 1.2, 0.06, 5.5, C_ACCENT)

# ════════════ SLIDE 13 — Testing & Deployment ════════════
slide13 = content_slide("Testing & Deployment")
# Testing
add_textbox(slide13, "Testing", 0.4, 1.3, 5.5, 0.4, font_size=15, bold=True, color=C_ACCENT2)
test_items = [
    "✅ Notes: CRUD verified",
    "✅ Flashcards: Manual + AI generation",
    "✅ Quiz: Generation + grading",
    "✅ Chat: Multi-turn AI conversations",
    "✅ UI/UX: Dark theme + responsiveness",
    "🔧 API: curl / Postman / DevTools",
]
add_bullet_box(slide13, test_items, 0.4, 1.8, 5.8, 4.0, font_size=13)

# Deployment
add_textbox(slide13, "Deployment", 7.0, 1.3, 5.5, 0.4, font_size=15, bold=True, color=C_ACCENT2)
dep_steps = [
    "1. Python 3.8+ & Node.js 18+ installed",
    "2. Get free Groq API key at console.groq.com",
    "3. cd backend → pip install -r requirements.txt",
    "4. Add GROQ_API_KEY to .env",
    "5. python run.py  (Port 5000)",
    "6. cd frontend → npm install → npm run dev  (Port 5173)",
]
add_bullet_box(slide13, dep_steps, 7.0, 1.8, 5.9, 4.0, font_size=13)

add_rect(slide13, 6.6, 1.2, 0.06, 5.5, C_ACCENT)

# ════════════ SLIDE 14 — Future & Conclusion ════════════
slide14 = content_slide("Future Enhancements & Conclusion")
future = [
    ("🔐 High",   "User Authentication",  "JWT-based multi-user login"),
    ("🗄  High",   "Database Migration",   "PostgreSQL / MongoDB"),
    ("📄 Medium", "File Upload & OCR",    "PDF/image text extraction"),
    ("🔁 Medium", "Spaced Repetition",    "SM-2 flashcard scheduling"),
    ("📱 Low",    "Mobile App",           "React Native companion"),
]
add_table_ppt(slide14,
    ["Priority", "Feature", "Description"],
    [(p, f, d) for p, f, d in future],
    0.4, 1.5, 12.5
)

# Conclusion box
add_rect(slide14, 0.4, 4.9, 12.5, 1.9, PptxRGB(0x1E, 0x1E, 0x44))
add_rect(slide14, 0.4, 4.9, 0.08, 1.9, C_ACCENT)
add_textbox(slide14,
    "Study Assistant AI successfully demonstrates how modern AI can be integrated into "
    "educational tools to create a powerful, personalized learning experience — combining "
    "React + Flask reliability with LLaMA 3.3 70B intelligence via Groq.",
    0.7, 5.0, 12.0, 1.7, font_size=14, color=C_LGRAY)

# ── Save PPTX ──
pptx_path = os.path.join(os.path.dirname(__file__), "Study_Assistant_AI_Presentation.pptx")
prs.save(pptx_path)
print(f"✅ PPTX saved: {pptx_path}")
